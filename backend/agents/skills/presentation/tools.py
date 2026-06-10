"""
agents/skills/presentation/tools.py
────────────────────────────────────
Core PPTX generation capability — shared by PresentationAgent (chatbot path)
and engagement_worker.py (async worker path).

Exported surface
────────────────
  PresentationResult                        — dataclass returned by all generators
  generate_from_template(name, tokens, id)  → PresentationResult
      Template path: downloads PPTX from blob-templates, substitutes {{tokens}},
      uploads personalised result to engagement-assets.  Used by the worker for
      all service onboarding decks.
  generate_presentation_core(context)       → PresentationResult
      LLM path: used by the chatbot (Ask AI) when a visitor requests a deck
      mid-conversation without having registered.

Note: this file intentionally does NOT export SCHEMAS or IMPLEMENTATIONS.
PresentationAgent overrides ask() directly — no tool-calling loop is used.

Environment variables (all injected by Terraform)
──────────────────────────────────────────────────
  AZURE_STORAGE_BLOB_ENDPOINT   — blob service endpoint
  ENGAGEMENT_ASSETS_CONTAINER   — output container (default: engagement-assets)
  BLOB_TEMPLATES_CONTAINER      — template source container (default: blob-templates)
  AZURE_OPENAI_KEY / ENDPOINT / MODEL / API_VERSION
"""

from __future__ import annotations

import io
import json
import logging
import os
import re
import uuid
import zipfile
from dataclasses import dataclass
from datetime import datetime, timedelta, timezone
from urllib.parse import urlparse

from azure.identity.aio import DefaultAzureCredential, ManagedIdentityCredential
from azure.storage.blob import BlobSasPermissions, generate_blob_sas
from azure.storage.blob.aio import BlobServiceClient
from openai import AzureOpenAI
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)

# ── Config ─────────────────────────────────────────────────────────────────────

_BLOB_ENDPOINT        = os.getenv("AZURE_STORAGE_BLOB_ENDPOINT", "")
_ASSETS_CONTAINER     = os.getenv("ENGAGEMENT_ASSETS_CONTAINER", "engagement-assets")
_TEMPLATES_CONTAINER  = os.getenv("BLOB_TEMPLATES_CONTAINER", "blob-templates")
_SAS_EXPIRY_HOURS     = 24

# ── Brand palette ──────────────────────────────────────────────────────────────

_TEAL      = RGBColor(0x0D, 0x94, 0x88)
_TEAL_DARK = RGBColor(0x0F, 0x76, 0x6E)
_DARK_BG   = RGBColor(0x0F, 0x0D, 0x18)
_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
_GRAY_TEXT = RGBColor(0x6B, 0x72, 0x80)
_BODY_TEXT = RGBColor(0x37, 0x41, 0x51)


# ── Result dataclass ───────────────────────────────────────────────────────────

@dataclass
class PresentationResult:
    pptx_bytes: bytes
    blob_path:  str
    sas_url:    str
    expires_at: datetime


# ── Credential helper ──────────────────────────────────────────────────────────

def _credential():
    client_id = os.getenv("AZURE_CLIENT_ID")
    if client_id:
        return ManagedIdentityCredential(client_id=client_id)
    return DefaultAzureCredential()


# ── LLM slide content generation ───────────────────────────────────────────────

_SLIDE_PROMPT = """\
You are a professional business writer for Terian Services, a technology consultancy
specialising in AI analytics, data engineering, and cloud solutions.

A prospect has expressed interest. Generate concise, professional slide content for a
6-slide onboarding deck tailored to their context.

Respond ONLY with a valid JSON object matching this exact schema (no markdown fences):
{{
  "slides": [
    {{
      "title": "string — slide heading (max 8 words)",
      "bullets": ["string", "string", "string"]
    }}
  ]
}}

Produce exactly 6 slide objects (indices 0–5):
  0 — Who We Are          (Terian Services overview, 2–3 lines of credibility)
  1 — Our Approach        (how Terian delivers engagements, methodology)
  2 — Your Engagement     (specific to engagement_type and tier, what they get)
  3 — Tailored for {org_name}  (industry fit, use_case alignment, user_count scale)
  4 — What Happens Next   (onboarding timeline: discovery → proposal → kickoff)
  5 — Let's Get Started   (contact info bullets, sales email, website, response time)

Context:
  Organization   : {org_name}
  Industry       : {industry}
  Estimated users: {user_count}
  Engagement type: {engagement_type}
  Tier           : {tier_interest}
  Use case       : {use_case}
"""

_FALLBACK_SLIDES = [ 
    {"title": "Who We Are",        "bullets": ["AI-native technology consultancy", "Deep expertise in data & cloud", "Proven delivery across industries"]},
    {"title": "Our Approach",      "bullets": ["Discovery → Design → Deliver", "Iterative, outcome-focused methodology", "Transparent progress at every stage"]},
    {"title": "Your Engagement",   "bullets": ["Engagement scoped to your needs", "Dedicated team from day one", "Clear milestones and deliverables"]},
    {"title": "Built for You",     "bullets": ["Industry-specific patterns applied", "Scales with your organisation", "Aligned to your use case"]},
    {"title": "What Happens Next", "bullets": ["Discovery call within 2 business days", "Tailored proposal within 1 week", "Kickoff on your timeline"]},
    {"title": "Let's Get Started", "bullets": ["sales@terian-services.com", "terian-services.com", "We respond within one business day"]},
]


def _call_llm_for_slides(context: dict) -> list[dict]:
    """Synchronous LLM call — run via asyncio.to_thread."""
    prompt = _SLIDE_PROMPT.format(
        org_name=context.get("org_name", "your organisation"),
        industry=context.get("industry") or "Technology",
        user_count=f"{context.get('user_count', 0):,}" if isinstance(context.get("user_count"), int) else str(context.get("user_count", "N/A")),
        engagement_type=context.get("engagement_type", ""),
        tier_interest=context.get("tier_interest", ""),
        use_case=context.get("use_case") or "Not specified",
    )
    client = AzureOpenAI(
        api_key=os.getenv("AZURE_OPENAI_KEY", ""),
        azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT", ""),
        api_version=os.getenv("AZURE_OPENAI_API_VERSION", "2024-12-01-preview"),
    )
    response = client.chat.completions.create(
        model=os.getenv("AZURE_OPENAI_MODEL", "gpt-4.1"),
        messages=[{"role": "user", "content": prompt}],
        max_completion_tokens=1200,
        temperature=0.4,
    )
    raw = (response.choices[0].message.content or "").strip()
    if raw.startswith("```"):
        raw = raw.split("```")[1]
        if raw.startswith("json"):
            raw = raw[4:]
    data = json.loads(raw)
    slides = data.get("slides", [])
    if len(slides) != 6:
        raise ValueError(f"Expected 6 slides from LLM, got {len(slides)}")
    return slides


# ── PPTX builder ───────────────────────────────────────────────────────────────

def _add_text_box(slide, left, top, width, height, text, font_size, color,
                  bold=False, align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p


def _build_pptx_bytes(context: dict, slides_content: list[dict]) -> bytes:
    """Build a Terian-branded .pptx deck and return raw bytes. Synchronous."""
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)

    W = prs.slide_width
    H = prs.slide_height
    HEADER_H = Inches(1.1)
    FOOTER_H = Inches(0.45)
    BODY_TOP  = HEADER_H
    MARGIN_L  = Inches(0.55)

    blank = prs.slide_layouts[6]

    # ── Cover slide ────────────────────────────────────────────────────────────
    cover = prs.slides.add_slide(blank)
    cover.background.fill.solid()
    cover.background.fill.fore_color.rgb = _DARK_BG

    accent = cover.shapes.add_shape(1, 0, 0, Inches(0.35), H)
    accent.fill.solid()
    accent.fill.fore_color.rgb = _TEAL
    accent.line.fill.background()

    _add_text_box(cover, MARGIN_L, Inches(0.35), W - Inches(0.7), Inches(0.5),
                  "TERIAN SERVICES", 11, _TEAL, bold=True)

    first_name = (context.get("full_name") or context.get("org_name") or "").split()[0]
    _add_text_box(cover, MARGIN_L, Inches(1.1), W - Inches(0.7), Inches(0.7),
                  f"Welcome, {first_name}" if first_name else "Welcome", 24, _WHITE)

    _add_text_box(cover, MARGIN_L, Inches(1.75), W - Inches(0.7), Inches(1.1),
                  context.get("org_name", ""), 36, _WHITE, bold=True)

    _add_text_box(cover, MARGIN_L, Inches(2.95), W - Inches(0.7), Inches(0.5),
                  "Your Onboarding Overview", 16, _TEAL)

    badge = f"{context.get('engagement_type', '')}  ·  {context.get('tier_interest', '')} Tier"
    _add_text_box(cover, MARGIN_L, Inches(3.6), W - Inches(0.7), Inches(0.4),
                  badge.strip(" · "), 13, _GRAY_TEXT)

    ftr = cover.shapes.add_shape(1, 0, H - FOOTER_H, W, FOOTER_H)
    ftr.fill.solid()
    ftr.fill.fore_color.rgb = RGBColor(0x1A, 0x17, 0x2B)
    ftr.line.fill.background()
    _add_text_box(cover, MARGIN_L, H - FOOTER_H + Pt(4), W - 2 * MARGIN_L, FOOTER_H,
                  "terian-services.com", 10, _GRAY_TEXT)

    # ── Content slides ─────────────────────────────────────────────────────────
    for slide_data in slides_content:
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _WHITE

        hdr = slide.shapes.add_shape(1, 0, 0, W, HEADER_H)
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = _TEAL
        hdr.line.fill.background()

        _add_text_box(slide, MARGIN_L, Inches(0.25), W - 2 * MARGIN_L, Inches(0.7),
                      slide_data.get("title", ""), 24, _WHITE, bold=True)

        bullet_top = BODY_TOP + Inches(0.35)
        for bullet in slide_data.get("bullets", [])[:4]:
            _add_text_box(slide, MARGIN_L, bullet_top,
                          Inches(0.25), Inches(0.4), "▸", 13, _TEAL, bold=True)
            _add_text_box(slide, MARGIN_L + Inches(0.3), bullet_top,
                          W - MARGIN_L - Inches(0.6), Inches(0.4),
                          bullet, 16, _BODY_TEXT)
            bullet_top += Inches(0.72)

        ftr2 = slide.shapes.add_shape(1, 0, H - FOOTER_H, W, FOOTER_H)
        ftr2.fill.solid()
        ftr2.fill.fore_color.rgb = _DARK_BG
        ftr2.line.fill.background()
        _add_text_box(slide, MARGIN_L, H - FOOTER_H + Pt(4),
                      W - 2 * MARGIN_L, FOOTER_H,
                      "Terian Services  ·  terian-services.com", 10, _GRAY_TEXT)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Blob upload + SAS URL ──────────────────────────────────────────────────────

async def _upload_and_sign(blob_path: str, pptx_bytes: bytes) -> tuple[str, datetime]:
    """
    Upload the PPTX to Blob Storage and return (sas_url, expiry).
    Uses a User Delegation SAS so no storage account key is needed —
    the UAMI must have the 'Storage Blob Delegator' role on the account.
    """
    expiry = datetime.now(timezone.utc) + timedelta(hours=_SAS_EXPIRY_HOURS)

    async with _credential() as cred:
        async with BlobServiceClient(account_url=_BLOB_ENDPOINT, credential=cred) as svc:
            # Upload
            blob = svc.get_blob_client(container=_ASSETS_CONTAINER, blob=blob_path)
            await blob.upload_blob(pptx_bytes, overwrite=True)

            # User delegation key for SAS signing
            delegation_key = await svc.get_user_delegation_key(
                key_start_time=datetime.now(timezone.utc),
                key_expiry_time=expiry,
            )

    account_name = urlparse(_BLOB_ENDPOINT).hostname.split(".")[0]
    sas_token = generate_blob_sas(
        account_name=account_name,
        container_name=_ASSETS_CONTAINER,
        blob_name=blob_path,
        user_delegation_key=delegation_key,
        permission=BlobSasPermissions(read=True),
        expiry=expiry,
    )
    sas_url = f"{_BLOB_ENDPOINT.rstrip('/')}/{_ASSETS_CONTAINER}/{blob_path}?{sas_token}"
    logger.info("presentation: uploaded → %s/%s (SAS valid %dh)", _ASSETS_CONTAINER, blob_path, _SAS_EXPIRY_HOURS)
    return sas_url, expiry


async def _upload_company_copy(company_copy_slug: str, pptx_bytes: bytes) -> None:
    """
    Upload an additional archival copy of a generated deck to a stable,
    company-named path under the "chatbot" folder of engagement-assets:

        chatbot/{company_copy_slug}-award-nomination-presentation.pptx

    This keeps one discoverable, overwrite-on-regenerate copy per
    organisation across all three generation entry points (engagement
    worker, chatbot, intro CTA). No SAS URL is generated — this copy is
    for internal/sales browsing, not for client delivery.

    Best-effort: failures are logged but never raised, so this never
    breaks the primary generation/delivery flow.
    """
    blob_path = f"chatbot/{company_copy_slug}-award-nomination-presentation.pptx"
    try:
        async with _credential() as cred:
            async with BlobServiceClient(account_url=_BLOB_ENDPOINT, credential=cred) as svc:
                blob = svc.get_blob_client(container=_ASSETS_CONTAINER, blob=blob_path)
                await blob.upload_blob(pptx_bytes, overwrite=True)
        logger.info("presentation: company copy saved → %s/%s", _ASSETS_CONTAINER, blob_path)
    except Exception as exc:
        logger.warning("presentation: failed to save company copy %s: %s", blob_path, exc)


# ── Template-based generation (worker path) ────────────────────────────────────

import asyncio  # noqa: E402 (placed here to keep top-of-file imports clean)


async def _download_template(template_name: str) -> bytes:
    """
    Download a PPTX template from the blob-templates container.
    Raises if the blob does not exist or the endpoint is not configured.
    """
    if not _BLOB_ENDPOINT:
        raise EnvironmentError("AZURE_STORAGE_BLOB_ENDPOINT is not set")
    async with _credential() as cred:
        async with BlobServiceClient(account_url=_BLOB_ENDPOINT, credential=cred) as svc:
            blob = svc.get_blob_client(container=_TEMPLATES_CONTAINER, blob=template_name)
            stream = await blob.download_blob()
            return await stream.readall()


def _substitute_tokens(pptx_bytes: bytes, tokens: dict[str, str]) -> bytes:
    """
    Replace {{TOKEN}} placeholders inside every slide XML in the PPTX zip.

    Operates at the zip entry level — reads each ppt/slides/slide*.xml,
    does a plain string replacement for every token, and writes the result
    back.  All other zip entries (images, layouts, theme, etc.) are copied
    unchanged so the visual design is fully preserved.

    This approach is safe for the award_nomination_onboarding template because
    all tokens were verified to appear as whole strings inside single <a:t>
    elements (no split-run issues).
    """
    src = io.BytesIO(pptx_bytes)
    dst = io.BytesIO()
    with zipfile.ZipFile(src, "r") as zin, \
         zipfile.ZipFile(dst, "w", compression=zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.startswith("ppt/slides/") and item.filename.endswith(".xml"):
                text = data.decode("utf-8")
                for key, value in tokens.items():
                    text = text.replace(f"{{{{{key}}}}}", value)
                data = text.encode("utf-8")
            zout.writestr(item, data)
    return dst.getvalue()


async def generate_from_template(
    template_name: str,
    tokens: dict[str, str],
    engagement_id: str | None = None,
    company_copy_slug: str | None = None,
) -> PresentationResult:
    """
    Template-based PPTX generation.

    1. Download the named template from blob-templates container.
    2. Substitute {{tokens}} across all slide XML (zip-level replace).
    3. Upload the personalised PPTX to engagement-assets.
    4. If company_copy_slug is set, also save an archival copy to
       chatbot/{company_copy_slug}-award-nomination-presentation.pptx
       (best-effort, does not affect the returned result).
    5. Return PresentationResult (pptx_bytes, blob_path, sas_url, expires_at).

    template_name      — blob name in blob-templates (e.g. "award_nomination_onboarding.pptx")
    tokens             — dict of TOKEN_KEY → replacement value (no braces in keys)
    engagement_id      — used as the blob path prefix; falls back to a random UUID prefix
    company_copy_slug  — if set, also archive a copy under chatbot/{slug}-award-nomination-presentation.pptx
    """
    # 1. Download template
    template_bytes = await _download_template(template_name)
    logger.info("presentation: downloaded template %s (%d bytes)", template_name, len(template_bytes))

    # 2. Substitute tokens in a thread (CPU-bound zip work)
    pptx_bytes = await asyncio.to_thread(_substitute_tokens, template_bytes, tokens)
    logger.info("presentation: token substitution complete (%d tokens)", len(tokens))

    # 3. Upload personalised deck + generate SAS URL
    blob_path = (
        f"{engagement_id}/onboarding.pptx"
        if engagement_id
        else f"chatbot/{uuid.uuid4()}/onboarding.pptx"
    )
    sas_url, expiry = await _upload_and_sign(blob_path, pptx_bytes)

    # 4. Archival company copy (best-effort, additional to the path above)
    if company_copy_slug:
        await _upload_company_copy(company_copy_slug, pptx_bytes)

    return PresentationResult(
        pptx_bytes=pptx_bytes,
        blob_path=blob_path,
        sas_url=sas_url,
        expires_at=expiry,
    )


# ── Product / services summary generation ──────────────────────────────────────

_PRODUCT_CONTENT_PROMPT = """\
You are a product content writer for Terian Services.

Generate concise, professional slide content for a {slide_count}-slide \
"{presentation_label}" presentation about {product_name}.

Respond ONLY with a valid JSON object (no markdown fences):
{{
  "slides": [
    {{
      "title": "string — slide heading (max 8 words)",
      "bullets": ["string", "string", "string"]
    }}
  ]
}}

Produce exactly {slide_count} slide objects matching these themes:
{slide_themes}

Product context:
  Product : {product_name}
  Tagline : {product_tagline}
  Overview: {product_overview}
"""

# Static product knowledge — source of truth until a dedicated product API exists.
_PRODUCTS: dict[str, dict] = {
    "award-nomination": {
        "name": "Award Nomination System",
        "tagline": "End-to-end digital nomination management for enterprise HR teams",
        "overview": (
            "A cloud-native SaaS platform that digitises and audits the full award "
            "nomination lifecycle — from employee submission through HRBP review, "
            "AI-assisted scoring, multi-tier approvals, and winner announcement. "
            "Built on Azure, deployed as a Container App, backed by Cosmos DB. "
            "Key capabilities: configurable nomination forms, graph-based fraud "
            "detection (ring/super-nominator/copy-paste/transactional patterns), "
            "real-time dashboards, SCIM provisioning, and audit trail export."
        ),
        "screen_flows_themes": (
            "0 — Employee Portal: submission flow (form → draft → submit)\n"
            "1 — HRBP Review Queue: scoring, filtering, bulk actions\n"
            "2 — Approval Chain: multi-tier approve / reject / escalate\n"
            "3 — Admin Dashboard: cycle management, analytics, exports\n"
            "4 — Notifications & Comms: email triggers at each stage\n"
            "5 — Winner Announcement: certificate generation and delivery"
        ),
        "infrastructure_themes": (
            "0 — High-Level Architecture: Azure Container Apps, Cosmos DB, Blob Storage\n"
            "1 — Authentication & Identity: Azure AD, SCIM, RBAC\n"
            "2 — Data Layer: Cosmos DB partitioning strategy, Blob containers\n"
            "3 — AI & Fraud Detection: OpenAI integration, graph analysis pipeline\n"
            "4 — DevOps & CI/CD: GitHub Actions, Terraform IaC, environments\n"
            "5 — Security & Compliance: Key Vault, TLS, audit logging, GDPR posture"
        ),
        "onboarding_themes": (
            "0 — What Is Award Nomination System: product overview and value proposition\n"
            "1 — Key Features: top 3-4 differentiating capabilities\n"
            "2 — How It Works: end-to-end nomination lifecycle in 4-5 steps\n"
            "3 — AI & Integrity Features: fraud detection, scoring, audit trail\n"
            "4 — Getting Started: onboarding steps, timelines, what to expect\n"
            "5 — Contact & Next Steps: how to reach sales, response time, resources"
        ),
    },
    "integrity-sentinel": {
        "name": "Integrity Sentinel",
        "tagline": "Real-time fraud and anomaly detection for enterprise business processes",
        "overview": (
            "An AI-powered integrity monitoring platform that runs continuously over "
            "HR, finance, and operational data to surface fraud, collusion, and policy "
            "violations before they become material losses. Uses a layered detection "
            "engine: rule-based thresholds, statistical anomaly detection, graph "
            "analysis (ring networks, super-nominators), and ML classification. "
            "Delivers findings via real-time dashboards, scored alerts, and "
            "investigator workflows — all deployable inside the client's Azure tenant."
        ),
        "screen_flows_themes": (
            "0 — Alert Feed: real-time scored alert list with severity filters\n"
            "1 — Case Investigation: drill-down view — evidence, graph, audit trail\n"
            "2 — Graph Explorer: interactive network visualisation of entity relationships\n"
            "3 — Rules & Thresholds: configuration UI for rule engine parameters\n"
            "4 — Analytics Dashboard: trend charts, detection rates, false-positive tuning\n"
            "5 — Reporting: scheduled reports, export to PDF/Excel, regulator packs"
        ),
        "infrastructure_themes": (
            "0 — High-Level Architecture: ingest pipeline, detection engine, alert store\n"
            "1 — Data Ingestion: connectors (HR, ERP, finance), batch + streaming\n"
            "2 — Detection Engine: rule layer → statistical layer → graph layer → ML layer\n"
            "3 — Azure Deployment: Container Apps, Cosmos DB, Event Hub, Key Vault\n"
            "4 — Tenant Isolation: deployed inside client Azure tenant, no data egress\n"
            "5 — Security & Compliance: encryption at rest/transit, RBAC, audit logging"
        ),
        "onboarding_themes": (
            "0 — What Is Integrity Sentinel: product overview and the problem it solves\n"
            "1 — Detection Capabilities: rule engine, statistical, graph, ML layers\n"
            "2 — How It Works: ingest → detect → alert → investigate → resolve\n"
            "3 — Key Differentiators: tenant isolation, explainability, low false-positive rate\n"
            "4 — Getting Started: data connectors, configuration, onboarding timeline\n"
            "5 — Contact & Next Steps: reach sales, pilot engagement offer, resources"
        ),
    },
}

_SERVICES_SLIDES: list[dict] = [
    {
        "title": "AI Analytics",
        "bullets": [
            "Custom ML models: forecasting, classification, anomaly detection, NLP",
            "Embedding-based search and GenAI pipelines with retrieval + evals/guardrails",
            "End-to-end delivery from problem framing through production deployment",
        ],
    },
    {
        "title": "Integrity & Fraud Detection",
        "bullets": [
            "Business-process-aware fraud detection across HR, finance, and operations",
            "Layered engine: rule thresholds → statistical anomaly → graph analysis → ML",
            "Patterns detected: collusion rings, duplicate payments, ghost vendors, copy-paste fraud",
        ],
    },
    {
        "title": "Data Mining",
        "bullets": [
            "Pattern discovery and driver analysis on operational, financial, and HR data",
            "Written findings reports paired with interactive dashboards (Power BI / Fabric / Looker)",
            "Delivered as a fixed-scope engagement with clear outputs and timelines",
        ],
    },
    {
        "title": "Datacenter → Cloud Migration",
        "bullets": [
            "Azure-first migration programs covering the 6 R's (rehost / replatform / refactor…)",
            "Structured approach: Assess → Design → Migrate → Optimize → Operate",
            "Delivered inside your Azure tenant under MSA & DPA — no data leaves your boundary",
        ],
    },
    {
        "title": "MLOps & Model Governance",
        "bullets": [
            "Model registry, drift monitoring, and automated evaluation harnesses",
            "Responsible-AI reviews, audits, and model cards for regulatory compliance",
            "Operates inside your Azure tenant for sensitive or regulated environments",
        ],
    },
    {
        "title": "Get In Touch",
        "bullets": [
            "sales@terian-services.com — we respond within one business day",
            "terian-services.com — full service descriptions and case studies",
            "All engagements scoped under MSA & DPA before any data is shared",
        ],
    },
]


def _call_llm_for_product_content(
    product_id: str,
    presentation_type: str,  # "onboarding" | "screen_flows" | "infrastructure"
) -> list[dict]:
    """Return 6 slides of structured content for the given product + deck type."""
    product = _PRODUCTS.get(product_id)
    if not product:
        raise ValueError(f"Unknown product_id: {product_id!r}")

    themes_key = f"{presentation_type}_themes"
    themes = product.get(themes_key, product["onboarding_themes"])

    label_map = {
        "onboarding":      "Product Onboarding",
        "screen_flows":    "Screen Flows",
        "infrastructure":  "Infrastructure Overview",
    }
    label = label_map.get(presentation_type, "Overview")

    prompt = _PRODUCT_CONTENT_PROMPT.format(
        slide_count=6,
        presentation_label=label,
        product_name=product["name"],
        product_tagline=product["tagline"],
        product_overview=product["overview"],
        slide_themes=themes,
    )

    client = AzureOpenAI(
        api_key=os.getenv("AZURE_OPENAI_KEY", ""),
        azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT", ""),
        api_version=os.getenv("AZURE_OPENAI_API_VERSION", "2024-12-01-preview"),
    )
    response = client.chat.completions.create(
        model=os.getenv("AZURE_OPENAI_MODEL", "gpt-4.1"),
        messages=[{"role": "user", "content": prompt}],
        max_completion_tokens=1400,
        temperature=0.3,
    )
    raw = (response.choices[0].message.content or "").strip()
    if raw.startswith("```"):
        raw = raw.split("```")[1]
        if raw.startswith("json"):
            raw = raw[4:]
    data = json.loads(raw)
    slides = data.get("slides", [])
    if len(slides) != 6:
        raise ValueError(f"Expected 6 slides, got {len(slides)}")
    return slides


def _build_summary_pptx_bytes(
    title: str,
    subtitle: str,
    slides_content: list[dict],
) -> bytes:
    """
    Build a clean Terian-branded PPTX from a title + list of {title, bullets} slides.
    The cover slide uses the dark brand background; content slides are white.
    """
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)

    W        = prs.slide_width
    H        = prs.slide_height
    HEADER_H = Inches(1.1)
    FOOTER_H = Inches(0.45)
    MARGIN_L = Inches(0.55)
    blank    = prs.slide_layouts[6]

    # ── Cover slide ────────────────────────────────────────────────────────────
    cover = prs.slides.add_slide(blank)
    cover.background.fill.solid()
    cover.background.fill.fore_color.rgb = _DARK_BG

    accent = cover.shapes.add_shape(1, 0, 0, Inches(0.35), H)
    accent.fill.solid()
    accent.fill.fore_color.rgb = _TEAL
    accent.line.fill.background()

    _add_text_box(cover, MARGIN_L, Inches(0.3), W - Inches(0.7), Inches(0.45),
                  "TERIAN SERVICES", 11, _TEAL, bold=True)
    _add_text_box(cover, MARGIN_L, Inches(1.2), W - Inches(0.7), Inches(1.1),
                  title, 34, _WHITE, bold=True)
    _add_text_box(cover, MARGIN_L, Inches(2.45), W - Inches(0.7), Inches(0.5),
                  subtitle, 16, _TEAL)

    ftr = cover.shapes.add_shape(1, 0, H - FOOTER_H, W, FOOTER_H)
    ftr.fill.solid()
    ftr.fill.fore_color.rgb = RGBColor(0x1A, 0x17, 0x2B)
    ftr.line.fill.background()
    _add_text_box(cover, MARGIN_L, H - FOOTER_H + Pt(4), W - 2 * MARGIN_L, FOOTER_H,
                  "terian-services.com", 10, _GRAY_TEXT)

    # ── Content slides ─────────────────────────────────────────────────────────
    for slide_data in slides_content:
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _WHITE

        hdr = slide.shapes.add_shape(1, 0, 0, W, HEADER_H)
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = _TEAL
        hdr.line.fill.background()

        _add_text_box(slide, MARGIN_L, Inches(0.25), W - 2 * MARGIN_L, Inches(0.7),
                      slide_data.get("title", ""), 24, _WHITE, bold=True)

        bullet_top = HEADER_H + Inches(0.35)
        for bullet in slide_data.get("bullets", [])[:4]:
            _add_text_box(slide, MARGIN_L, bullet_top,
                          Inches(0.25), Inches(0.4), "▸", 13, _TEAL, bold=True)
            _add_text_box(slide, MARGIN_L + Inches(0.3), bullet_top,
                          W - MARGIN_L - Inches(0.6), Inches(0.4),
                          bullet, 15, _BODY_TEXT)
            bullet_top += Inches(0.72)

        ftr2 = slide.shapes.add_shape(1, 0, H - FOOTER_H, W, FOOTER_H)
        ftr2.fill.solid()
        ftr2.fill.fore_color.rgb = _DARK_BG
        ftr2.line.fill.background()
        _add_text_box(slide, MARGIN_L, H - FOOTER_H + Pt(4),
                      W - 2 * MARGIN_L, FOOTER_H,
                      "Terian Services  ·  terian-services.com", 10, _GRAY_TEXT)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


async def generate_summary_presentation(
    product_id: str,
    presentation_type: str,
    blob_prefix: str = "chatbot",
) -> PresentationResult:
    """
    LLM → PPTX → Blob pipeline for product summary presentations.

    product_id        — "award-nomination" | "integrity-sentinel"
    presentation_type — "onboarding" | "screen_flows" | "infrastructure"
    blob_prefix       — path prefix in engagement-assets container
    """
    if not _BLOB_ENDPOINT:
        raise EnvironmentError("AZURE_STORAGE_BLOB_ENDPOINT is not set")

    product = _PRODUCTS.get(product_id)
    if not product:
        raise ValueError(f"Unknown product_id: {product_id!r}")

    label_map = {
        "onboarding":     "Product Onboarding",
        "screen_flows":   "Screen Flows",
        "infrastructure": "Infrastructure Overview",
    }
    subtitle = label_map.get(presentation_type, "Overview")

    try:
        slides_content = await asyncio.to_thread(
            _call_llm_for_product_content, product_id, presentation_type
        )
    except Exception as exc:
        logger.warning(
            "presentation: LLM product content failed (%s) — using fallback", exc
        )
        slides_content = _FALLBACK_SLIDES

    pptx_bytes = await asyncio.to_thread(
        _build_summary_pptx_bytes,
        product["name"],
        subtitle,
        slides_content,
    )

    blob_path = f"{blob_prefix}/{uuid.uuid4()}/{product_id}-{presentation_type}.pptx"
    sas_url, expiry = await _upload_and_sign(blob_path, pptx_bytes)

    logger.info(
        "presentation: generated %s/%s → %s", product_id, presentation_type, blob_path
    )
    return PresentationResult(
        pptx_bytes=pptx_bytes,
        blob_path=blob_path,
        sas_url=sas_url,
        expires_at=expiry,
    )


async def generate_services_overview_pptx(
    blob_prefix: str = "chatbot",
) -> PresentationResult:
    """
    Build a Terian Services overview presentation from the static services content.
    No LLM call — content is authoritative and hardcoded.
    """
    if not _BLOB_ENDPOINT:
        raise EnvironmentError("AZURE_STORAGE_BLOB_ENDPOINT is not set")

    pptx_bytes = await asyncio.to_thread(
        _build_summary_pptx_bytes,
        "Terian Services",
        "Our Engineering & Analytics Capabilities",
        _SERVICES_SLIDES,
    )

    blob_path = f"{blob_prefix}/{uuid.uuid4()}/services-overview.pptx"
    sas_url, expiry = await _upload_and_sign(blob_path, pptx_bytes)

    logger.info("presentation: generated services overview → %s", blob_path)
    return PresentationResult(
        pptx_bytes=pptx_bytes,
        blob_path=blob_path,
        sas_url=sas_url,
        expires_at=expiry,
    )


# ── LLM generation (chatbot / Ask AI path) ────────────────────────────────────

async def generate_presentation_core(context: dict) -> PresentationResult:
    """
    Full generation pipeline — used by both PresentationAgent.generate()
    (worker path) and the LLM tool wrapper (chatbot path).

    context keys (all optional — defaults applied if missing):
      org_name, full_name, industry, user_count,
      engagement_type, tier_interest, use_case
    """
    if not _BLOB_ENDPOINT:
        raise EnvironmentError("AZURE_STORAGE_BLOB_ENDPOINT is not set")

    # 1. LLM → slide content
    try:
        slides_content = await asyncio.to_thread(_call_llm_for_slides, context)
    except Exception as exc:
        logger.warning("presentation: LLM slide generation failed (%s) — using fallback", exc)
        slides_content = _FALLBACK_SLIDES

    # 2. Build PPTX in memory
    pptx_bytes = await asyncio.to_thread(_build_pptx_bytes, context, slides_content)

    # 3. Upload + sign
    blob_path = f"chatbot/{uuid.uuid4()}/onboarding.pptx"
    # Worker overrides blob_path via engagement_id — handled in worker.py
    if context.get("engagement_id"):
        blob_path = f"{context['engagement_id']}/onboarding.pptx"

    sas_url, expiry = await _upload_and_sign(blob_path, pptx_bytes)

    return PresentationResult(
        pptx_bytes=pptx_bytes,
        blob_path=blob_path,
        sas_url=sas_url,
        expires_at=expiry,
    )

